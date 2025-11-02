#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Microsoft Office ファイルを Markdown に変換するツール
サポートファイル: .docx, .pptx, .xlsx, .pdf
"""

import os
import sys
from pathlib import Path
from typing import Optional

# Word (.docx) 変換用
import mammoth

# PowerPoint (.pptx) 変換用
from pptx import Presentation

# Excel (.xlsx) 変換用
import openpyxl

# PDF 変換用
from pdfminer.high_level import extract_text as extract_pdf_text
from pdfminer.layout import LAParams

# Markdown 整形用
from markdownify import markdownify
import re


class MarkdownFormatter:
    """Markdownテキストを整形するクラス"""

    @staticmethod
    def remove_single_char_spaces(text: str) -> str:
        """1文字ずつスペースで区切られた文字列を結合する"""
        # パターン: 1-2文字 + スペース + 1-2文字 の繰り返し
        # 例: "A I " -> "AI", "自 己 紹 介" -> "自己紹介"
        pattern = r'(\S)\s+(?=\S\s|\S$)'

        # 行ごとに処理
        lines = text.split('\n')
        result_lines = []

        for line in lines:
            # スペース区切りの文字数をカウント
            words = line.split()
            if len(words) >= 3:  # 3つ以上のトークンがある場合
                # すべてのトークンが1-2文字かチェック
                if all(len(word) <= 2 for word in words):
                    # スペースを除去して結合
                    line = ''.join(words)
            result_lines.append(line)

        return '\n'.join(result_lines)

    @staticmethod
    def format(markdown_text: str) -> str:
        """Markdownテキストを整形する"""
        # まずスペース区切りの文字を結合
        markdown_text = MarkdownFormatter.remove_single_char_spaces(markdown_text)

        lines = markdown_text.split('\n')
        formatted_lines = []
        buffer = []

        for i, line in enumerate(lines):
            stripped = line.strip()

            # 1-2文字の行を一時バッファに蓄積
            if len(stripped) <= 2 and stripped and not stripped.startswith('#'):
                buffer.append(stripped)
            else:
                # バッファに溜まった短い行を結合
                if buffer:
                    combined = ''.join(buffer)
                    # 結合した結果が意味のある長さなら追加
                    if len(combined) > 2:
                        formatted_lines.append(combined)
                    buffer = []

                # 現在の行を追加（空行でない場合）
                if stripped:
                    formatted_lines.append(line.rstrip())
                elif formatted_lines and formatted_lines[-1] != '':
                    # 空行は1つだけ保持
                    formatted_lines.append('')

        # 残ったバッファを処理
        if buffer:
            combined = ''.join(buffer)
            if len(combined) > 2:
                formatted_lines.append(combined)

        # 連続する空行を削減（最大2行まで）
        result = []
        empty_count = 0
        for line in formatted_lines:
            if line == '':
                empty_count += 1
                if empty_count <= 2:
                    result.append(line)
            else:
                empty_count = 0
                result.append(line)

        # 先頭と末尾の余分な空行を削除
        while result and result[0] == '':
            result.pop(0)
        while result and result[-1] == '':
            result.pop()

        return '\n'.join(result)


class FileConverter:
    """ファイルをMarkdownに変換するクラス"""

    def __init__(self, input_dir: str = "./input", output_dir: str = "./output"):
        self.input_dir = Path(input_dir)
        self.output_dir = Path(output_dir)

        # ディレクトリが存在することを確認
        if not self.input_dir.exists():
            raise FileNotFoundError(f"入力ディレクトリが見つかりません: {self.input_dir}")

        # 出力ディレクトリを作成
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def convert_docx(self, file_path: Path) -> str:
        """Word文書(.docx)をMarkdownに変換"""
        with open(file_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html = result.value
            # HTMLをMarkdownに変換
            markdown = markdownify(html)
            # 整形して返す
            return MarkdownFormatter.format(markdown)

    def convert_pptx(self, file_path: Path) -> str:
        """PowerPoint(.pptx)をMarkdownに変換"""
        prs = Presentation(file_path)
        markdown_lines = [f"# {file_path.stem}\n"]

        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_lines.append(f"\n## スライド {slide_num}\n")

            # テキストフレームからテキストを抽出
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        markdown_lines.append(f"{text}\n")

        markdown = "\n".join(markdown_lines)
        # 整形して返す
        return MarkdownFormatter.format(markdown)

    def convert_xlsx(self, file_path: Path) -> str:
        """Excel(.xlsx)をMarkdownに変換"""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        markdown_lines = [f"# {file_path.stem}\n"]

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            markdown_lines.append(f"\n## シート: {sheet_name}\n")

            # テーブル形式で出力
            rows = []
            for row in sheet.iter_rows(values_only=True):
                # 空行をスキップ
                if all(cell is None or str(cell).strip() == "" for cell in row):
                    continue
                # セルの値を文字列に変換
                row_values = [str(cell) if cell is not None else "" for cell in row]
                rows.append(row_values)

            if rows:
                # Markdownテーブルを作成
                if len(rows) > 0:
                    # ヘッダー行
                    header = "| " + " | ".join(rows[0]) + " |"
                    separator = "| " + " | ".join(["---"] * len(rows[0])) + " |"
                    markdown_lines.append(header)
                    markdown_lines.append(separator)

                    # データ行
                    for row in rows[1:]:
                        markdown_lines.append("| " + " | ".join(row) + " |")

                    markdown_lines.append("")

        markdown = "\n".join(markdown_lines)
        # 整形して返す（テーブルは保持）
        return MarkdownFormatter.format(markdown)

    def convert_pdf(self, file_path: Path) -> str:
        """PDF(.pdf)をMarkdownに変換"""
        # PDFからテキストを抽出
        laparams = LAParams()
        text = extract_pdf_text(str(file_path), laparams=laparams)

        # 基本的なMarkdown形式に整形
        markdown_lines = [f"# {file_path.stem}\n"]
        markdown_lines.append(text)

        markdown = "\n".join(markdown_lines)
        # 整形して返す
        return MarkdownFormatter.format(markdown)

    def convert_file(self, file_path: Path) -> Optional[str]:
        """ファイルの拡張子に応じて適切な変換を実行"""
        extension = file_path.suffix.lower()

        converters = {
            '.docx': self.convert_docx,
            '.pptx': self.convert_pptx,
            '.xlsx': self.convert_xlsx,
            '.pdf': self.convert_pdf,
        }

        converter = converters.get(extension)
        if converter:
            try:
                print(f"変換中: {file_path.name} ...", end=" ")
                markdown_content = converter(file_path)
                print("完了")
                return markdown_content
            except Exception as e:
                print(f"エラー: {str(e)}")
                return None
        else:
            print(f"スキップ: {file_path.name} (サポートされていない形式)")
            return None

    def convert_all(self):
        """入力ディレクトリ内のすべてのファイルを変換"""
        print(f"\n入力ディレクトリ: {self.input_dir}")
        print(f"出力ディレクトリ: {self.output_dir}\n")

        # サポートされている拡張子
        supported_extensions = {'.docx', '.pptx', '.xlsx', '.pdf'}

        # 入力ディレクトリ内のファイルを取得
        files = [f for f in self.input_dir.iterdir() if f.is_file() and f.suffix.lower() in supported_extensions]

        if not files:
            print("変換可能なファイルが見つかりません。")
            print(f"サポートされている形式: {', '.join(supported_extensions)}")
            return

        print(f"{len(files)}個のファイルを検出しました。\n")

        converted_count = 0
        for file_path in files:
            markdown_content = self.convert_file(file_path)

            if markdown_content:
                # 出力ファイル名を作成
                output_file = self.output_dir / f"{file_path.stem}.md"

                # Markdownファイルとして保存
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(markdown_content)

                converted_count += 1

        print(f"\n変換完了: {converted_count}/{len(files)} ファイル")
        print(f"出力先: {self.output_dir}")


def main():
    """メイン関数"""
    # デフォルトのディレクトリパス
    input_dir = "./input"
    output_dir = "./output"

    # コマンドライン引数からディレクトリパスを取得（オプション）
    if len(sys.argv) > 1:
        input_dir = sys.argv[1]
    if len(sys.argv) > 2:
        output_dir = sys.argv[2]

    try:
        converter = FileConverter(input_dir, output_dir)
        converter.convert_all()
    except Exception as e:
        print(f"エラーが発生しました: {str(e)}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
